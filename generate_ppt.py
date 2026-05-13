#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate 15-minute PPT for MATLAB-gprMax joint simulation report."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(BASE_DIR, 'ppt_images')
OUTPUT_PPT = os.path.join(BASE_DIR, 'MATLAB_gprMax_Simulation_Report.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_TITLE = RGBColor(0x1F, 0x4E, 0x78)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_TITLE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.0))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Engineering Ethics Course Report | 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets, image_name=None, image_left=None,
                      image_top=None, image_width=None, image_height=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_TITLE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if image_name and os.path.exists(os.path.join(IMG_DIR, image_name)):
        left_margin = Inches(0.5)
        text_width = Inches(5.5)
        if image_left is None:
            image_left = Inches(6.5)
        if image_top is None:
            image_top = Inches(1.4)
        if image_width is None:
            image_width = Inches(6.0)
        if image_height is None:
            image_height = Inches(5.5)
    else:
        left_margin = Inches(0.8)
        text_width = Inches(11.5)

    text_box = slide.shapes.add_textbox(left_margin, Inches(1.3), text_width, Inches(5.8))
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(12)
        p.level = 0

    if image_name and os.path.exists(os.path.join(IMG_DIR, image_name)):
        slide.shapes.add_picture(os.path.join(IMG_DIR, image_name),
                                 image_left, image_top,
                                 width=image_width, height=image_height)

    add_page_number(slide, len(prs.slides))
    return slide


def add_full_image_slide(prs, title, image_name):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_TITLE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if os.path.exists(os.path.join(IMG_DIR, image_name)):
        slide.shapes.add_picture(os.path.join(IMG_DIR, image_name),
                                 Inches(0.3), Inches(1.0),
                                 width=Inches(12.7), height=Inches(6.2))

    add_page_number(slide, len(prs.slides))
    return slide


def add_page_number(slide, number):
    box = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p.alignment = PP_ALIGN.RIGHT


add_title_slide(prs,
    "MATLAB + gprMax Joint Simulation for GPR Forward Modeling",
    "Complete Workflow and Engineering Ethics Discussion\n\nBased on FDTD method, HDF5 I/O, and MATLAB signal processing")

add_content_slide(prs, "Agenda",
    [
        "1.  gprMax Overview -- FDTD solver for Maxwell's equations",
        "2.  MATLAB's Role -- Parameter control and post-processing",
        "3.  Joint Simulation Workflow -- From model to visualization",
        "4.  Step 1: Input File (.in) Syntax and Modeling",
        "5.  Step 2: gprMax Forward Computation",
        "6.  Step 3: HDF5 Output File Structure",
        "7.  Step 4: MATLAB Data Reading (h5read)",
        "8.  Step 5: MATLAB Post-processing and Visualization",
        "9.  Step 6: Engineering Application and Format Conversion",
        "10. Summary and Engineering Ethics Considerations",
    ])

add_content_slide(prs, "gprMax: FDTD Electromagnetic Solver",
    [
        "gprMax solves Maxwell's curl equations using the Finite-Difference Time-Domain (FDTD) method",
        "Yee-cell discretization: spatial steps dx, dy, dz define the grid resolution",
        "Spatial resolution limits maximum frequency: typically <= lambda/10",
        "Supports PML absorbing boundaries, dispersive materials, real antenna models",
        "Parallel acceleration: OpenMP (CPU), MPI (distributed), CUDA (GPU)",
        "Output format: HDF5 with time-history field data from all receivers",
        "Originally designed for GPR, applicable to general EM simulation problems",
    ])

add_content_slide(prs, "MATLAB's Role in Joint Simulation",
    [
        "[Front-end: Parameterized Modeling] MATLAB batch-generates gprMax input files (.in)",
        "[Middle: Flow Control] MATLAB system() calls gprMax CLI for automated B-scans",
        "[Back-end: Data I/O] Built-in h5read / h5readatt directly read HDF5 outputs",
        "[Back-end: Visualization] A-scan / B-scan plotting, spectral analysis, time-frequency",
        "[Back-end: Format Conversion] Convert to commercial GPR formats (RD3/DZT/DT1)",
        "Advantage: MATLAB matrix operations and Signal Processing Toolbox fit GPR data perfectly",
    ])

add_full_image_slide(prs, "Joint Simulation Workflow", "fig_workflow.png")

add_full_image_slide(prs, "Step 1: gprMax Input File (.in)", "fig_input_file.png")

add_content_slide(prs, "Step 2: gprMax Forward Computation",
    [
        "Command: python -m gprMax input_file.in -n 60",
        "-n 60: Run 60 traces for B-scan, antenna position shifts per src_steps/rx_steps",
        "--geometry-fixed: Build geometry once, only move Tx/Rx afterwards (much faster)",
        "-gpu: Enable CUDA GPU acceleration for large-scale models",
        "Each run produces one .out file (HDF5 format)",
        "Multi-trace B-scans need outputfiles_merge to create merged.out",
    ],
    image_name="fig_input_file.png",
    image_left=Inches(7.5), image_top=Inches(1.3),
    image_width=Inches(5.5), image_height=Inches(5.5))

add_full_image_slide(prs, "Step 3: HDF5 Output File Structure", "fig_hdf5_structure.png")

add_full_image_slide(prs, "Step 4: MATLAB Reading gprMax Output", "fig_matlab_code.png")

add_full_image_slide(prs, "Step 5: A-scan Visualization -- 6 EM Field Components", "fig_Ascan.png")

add_full_image_slide(prs, "Step 5: B-scan Visualization -- Metal Cylinder Hyperbola", "fig_Bscan.png")

add_full_image_slide(prs, "Step 5: Spectrum Analysis -- Validate Center Frequency", "fig_spectrum.png")

add_content_slide(prs, "Step 6: Engineering Application -- Commercial GPR Format Conversion",
    [
        "MATLAB script outputfile_converter.m converts gprMax output to industry standards:",
        "  * RD3 (Mala GeoScience) -- 16-bit binary + RAD header file",
        "  * DZT (GSSI) -- Single binary file with embedded header",
        "  * DT1 (Sensors & Software) -- HD header + DT1 data file",
        "  * IPRB (Impulse Radar) -- IPRH header + IPRB data file",
        "Significance: Simulated data can be loaded into commercial GPR software (e.g., ReflexW)",
        "Engineering value: Direct comparison with field data on the same platform",
    ])

add_full_image_slide(prs, "Engineering Ethics: Responsibility of Using Simulation Results", "fig_ethics.png")

add_content_slide(prs, "Summary",
    [
        "MATLAB + gprMax forms a complete GPR forward simulation pipeline:",
        "  Modeling -> Solving -> Output -> Reading -> Visualization -> Conversion",
        "MATLAB excels at parameterized batch processing and signal processing",
        "gprMax excels at physically rigorous FDTD solving and flexible model definition",
        "Engineering ethics key points:",
        "  * Clarify model assumptions; do not equate ideal simulation with field measurement",
        "  * Parameter choices must have physical basis (wavelength sampling, PML distance, time window)",
        "  * Simulation aids decision-making but does not replace engineering judgment and safety verification",
    ])

prs.save(OUTPUT_PPT)
print(f"PPT saved to: {OUTPUT_PPT}")
print(f"Total slides: {len(prs.slides)}")
